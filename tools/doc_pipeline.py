#!/usr/bin/env python3
"""社内知恵袋: /doc コマンドのコード制御パイプライン（v1.0.68）

設計方針:
- モデルの役割は「目次の JSON 出力（構造化出力で強制）」と「セクションの文章化」のみ。
  判断・検索・ファイル化・最終応答はすべてこのモジュール（コード）が実行する。
- v1.0.68 の本質対応: 最終応答に LLM を使わない。
  生成した案内文を OpenAI 互換の SSE チャンク（StreamingResponse）に変換して
  標準の応答処理 (process_chat_response) に渡すため、UI 配信・DB 保存・
  タイトル生成など下流の全処理が通常チャットと同じ経路で動作する。
  3B モデルの復唱省略・脚色という問題を構造的に解消した。
- PDF リンクは delta.annotations (url_citation) 経由で「ソース」としても表示する。
- 生成 PDF は ShineosMcpoFiles (9003) の配信ディレクトリ (data/mcpo_output) に
  直接書き出し、既存の /files 配信でオフライン閲覧できる。古い PDF は 7 日で自動削除。

使い方: main.py の process_chat から
  maybe_build_doc_response(request, form_data, user, metadata, model)
を呼び、None 以外ならその戻り値を process_chat の戻り値として返す。
"""
import asyncio
import datetime
import glob
import html
import json
import os
import re
import time
from pathlib import Path

import httpx
from starlette.responses import StreamingResponse

DOC_TRIGGER = "/doc"
# コマンドごとに出力形式を決める（/doc は PDF 既定の別名）
TRIGGERS = {"/docx": "docx", "/pptx": "pptx", "/pdf": "pdf", "/doc": "pdf"}
MAX_SECTIONS = 5
OUTLINE_TIMEOUT = 60.0
SECTION_TIMEOUT = 120.0
PDF_RETENTION_DAYS = 7
REFUSAL_MARKER = "対象外"
SUPPORTED_FORMATS = ("pdf", "pptx", "docx")

# 日本語フォント（tools/filegen_server.py と同じ候補）
_FONT_CANDIDATES = [
    (r"C:\Windows\Fonts\meiryo.ttc", 0, "Meiryo"),
    (r"C:\Windows\Fonts\msgothic.ttc", 0, "MSGothic"),
    (r"C:\Windows\Fonts\yugothic.ttc", 0, "YuGothic"),
]


def _data_dir() -> Path:
    try:
        from open_webui.env import DATA_DIR

        return Path(DATA_DIR)
    except Exception:
        return Path(os.environ.get("DATA_DIR", r"C:\Program Files\ShineosQA\data"))


def _output_dir() -> Path:
    d = Path(os.environ.get("FILE_EXPORT_DIR") or (_data_dir() / "mcpo_output"))
    d.mkdir(parents=True, exist_ok=True)
    return d


def _file_base_url() -> str:
    return os.environ.get("FILE_EXPORT_BASE_URL") or "http://localhost:9003/files"


def _ollama_base(request) -> str:
    try:
        urls = getattr(request.app.state.config, "OLLAMA_BASE_URLS", None)
        if urls:
            return urls[0].rstrip("/")
    except Exception:
        pass
    return os.environ.get("OLLAMA_BASE_URL", "http://127.0.0.1:11434").rstrip("/")


def _base_model_id(request, model_id: str) -> str:
    try:
        m = request.app.state.MODELS.get(model_id) or {}
        mid = (m.get("info") or {}).get("base_model_id") or ""
        if mid:
            return mid
    except Exception:
        pass
    return model_id or "qwen2.5:3b"


def _cleanup_old_files(days: int = PDF_RETENTION_DAYS):
    """保存期間を過ぎた生成資料（pdf/pptx/docx）を削除する（ストレージ蓄積の防止）"""
    try:
        cutoff = time.time() - days * 86400
        for ext in ("pdf", "pptx", "docx"):
            for f in glob.glob(str(_output_dir() / f"doc_*.{ext}")):
                try:
                    if os.path.getmtime(f) < cutoff:
                        os.remove(f)
                except Exception:
                    continue
    except Exception:
        pass


async def _emit(emitter, description: str, done: bool = False):
    if not emitter:
        return
    try:
        await emitter({"type": "status", "data": {"action": "document", "description": description, "done": done}})
    except Exception:
        pass


async def _make_outline(request, model_id: str, topic: str) -> dict:
    """目次を Ollama の構造化出力（format=schema）で生成する。

    チャット用モデル（拒否システムプロンプト付き）を使わず素のベースモデルに
    直接問い合わせるため、Q&A ルールと干渉しない。
    """
    schema = {
        "type": "object",
        "properties": {
            "title": {"type": "string"},
            "sections": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "query": {"type": "string"},
                    },
                    "required": ["heading", "query"],
                },
            },
        },
        "required": ["title", "sections"],
    }
    prompt = (
        f"社内文書を検索して「{topic}」に関する資料を作ります。\n"
        "資料の表題と、各セクションの見出し・検索語を JSON で出力してください。\n"
        "セクションは2〜5個。検索語は社内文書を探すための短いキーワードにしてください。\n"
        '例: {"title": "出張手当の整理", "sections": [{"heading": "概要", "query": "出張手当"}]}'
    )
    body = {
        "model": _base_model_id(request, model_id),
        "messages": [{"role": "user", "content": prompt}],
        "stream": False,
        "format": schema,
        "options": {"temperature": 0.2, "num_predict": 400},
    }
    async with httpx.AsyncClient(timeout=OUTLINE_TIMEOUT) as client:
        r = await client.post(_ollama_base(request) + "/api/chat", json=body)
        r.raise_for_status()
        content = (r.json().get("message") or {}).get("content", "")
    try:
        data = json.loads(content)
        title = str(data.get("title") or topic)[:60]
        sections = [
            {"heading": str(s.get("heading") or "")[:60], "query": str(s.get("query") or s.get("heading") or "")[:60]}
            for s in (data.get("sections") or [])
            if isinstance(s, dict) and s.get("heading")
        ]
        sections = sections[:MAX_SECTIONS]
        if sections:
            return {"title": title, "sections": sections}
    except Exception:
        pass
    # フォールバック: 固定テンプレ
    return {
        "title": topic[:60],
        "sections": [
            {"heading": "概要", "query": topic[:40]},
            {"heading": "詳細", "query": topic[:40] + " 手順"},
            {"heading": "まとめ", "query": topic[:40]},
        ],
    }


_SECTION_PROMPT = (
    "社内文書の参考片段を使って、資料のセクション「{heading}」の本文を書いてください。\n"
    "ルール:\n"
    "- 与えられた社内文書の片段だけを使うこと（片段に無い事実・数値を書かない）\n"
    "- 箇条書き中心・簡潔にまとめること\n"
    "- 参考片段の中に記載が無い場合は「この質問は社内文書の対象外です。一般的な情報はチャット入力欄のWeb検索をONにすると回答できます。」とだけ答えること\n"
    "\n参考片段:\n{query}"
)


async def _write_section(request, form_data: dict, heading: str, query: str) -> str:
    """セクション本文を 1 回のチャット完結で書かせる。

    ナレッジ RAG 注入（legacy）はチャット経路に組み込み済みのため、
    ループバックの chat completions を 1 回呼ぶだけで「検索 + 文章化」が完結する。
    材料が無い場合は Q&A の拒否文が返るので、資料向けの表記に置き換える。
    """
    auth = request.headers.get("authorization") or ""
    headers = {"Authorization": auth, "Content-Type": "application/json"} if auth else {"Content-Type": "application/json"}
    payload = {
        "model": form_data.get("model"),
        "messages": [{"role": "user", "content": _SECTION_PROMPT.format(heading=heading, query=query)}],
        "stream": False,
    }
    base = str(request.base_url).rstrip("/")
    async with httpx.AsyncClient(timeout=SECTION_TIMEOUT) as client:
        r = await client.post(base + "/api/chat/completions", json=payload, headers=headers)
        r.raise_for_status()
        data = r.json()
    content = ((data.get("choices") or [{}])[0].get("message") or {}).get("content", "")
    content = str(content).strip()
    if not content or REFUSAL_MARKER in content:
        return "（社内資料に該当記載がありません）"
    # 引用表記は資料ではノイズになるため簡易に落とす
    content = re.sub(r"\[\d+\]", "", content)
    return content


def _render_pdf(title: str, sections: list[dict], out_path: Path):
    """節構造（目次つき）を PDF にする（reportlab・日本語フォント対応）。"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    font_name = None
    for fpath, sub, name in _FONT_CANDIDATES:
        if os.path.exists(fpath):
            try:
                pdfmetrics.registerFont(TTFont(name, fpath, subfontIndex=sub))
                font_name = name
                break
            except Exception:
                continue

    doc = SimpleDocTemplate(str(out_path), pagesize=A4, title=title)
    styles = getSampleStyleSheet()
    for s in styles.byName.values():
        try:
            if font_name:
                s.fontName = font_name
        except Exception:
            pass

    story = [Paragraph(html.escape(title), styles["Title"]), Spacer(1, 12)]
    # 目次
    story.append(Paragraph("目次", styles["Heading2"]))
    for idx, sec in enumerate(sections, 1):
        story.append(Paragraph(html.escape(f"{idx}. {sec['heading']}"), styles["BodyText"]))
    story.append(Spacer(1, 14))
    for idx, sec in enumerate(sections, 1):
        story.append(Paragraph(html.escape(f"{idx}. {sec['heading']}"), styles["Heading2"]))
        story.append(Spacer(1, 6))
        for line in sec["body"].splitlines():
            line = line.strip()
            if not line:
                continue
            if line.startswith(("-", "・", "•")):
                story.append(Paragraph("• " + html.escape(line.lstrip("-・• ").strip()), styles["BodyText"]))
            else:
                story.append(Paragraph(html.escape(line), styles["BodyText"]))
        story.append(Spacer(1, 10))
    doc.build(story)


def _build_document(fmt: str, title: str, sections: list[dict]) -> str:
    """節構造を指定形式（pdf/pptx/docx）で書き出し、配信 URL を返す"""
    _cleanup_old_files()
    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    fname = f"doc_{stamp}.{fmt}"
    out = _output_dir() / fname
    if fmt == "pdf":
        _render_pdf(title, sections, out)
    elif fmt == "pptx":
        _render_pptx(title, sections, out)
    else:
        _render_docx(title, sections, out)
    return f"{_file_base_url()}/{fname}"


def _render_pptx(title: str, sections: list[dict], out_path: Path):
    """スライド形式: 表題スライド + セクションごとに 1 枚（箇条書き）"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title

    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sec["heading"]
        body = slide.placeholders[1].text_frame
        body.text = ""
        first = True
        for line in sec["body"].splitlines():
            line = line.strip()
            if not line:
                continue
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            p.text = ("・" + line.lstrip("-・• ").strip()) if not line.startswith(("（", "(")) else line
            p.font.size = Pt(18)
    prs.save(str(out_path))


def _render_docx(title: str, sections: list[dict], out_path: Path):
    """Word 形式: 見出し + 本文（PDF と同じ節構造）"""
    import docx

    d = docx.Document()
    d.add_heading(title, 0)
    d.add_heading("目次", level=1)
    for idx, sec in enumerate(sections, 1):
        d.add_paragraph(f"{idx}. {sec['heading']}")
    for idx, sec in enumerate(sections, 1):
        d.add_heading(f"{idx}. {sec['heading']}", level=1)
        for line in sec["body"].splitlines():
            line = line.strip()
            if not line:
                continue
            if line.startswith(("-", "・", "•")):
                d.add_paragraph(line.lstrip("-・• ").strip(), style="List Bullet")
            else:
                d.add_paragraph(line)
    d.save(str(out_path))


def _sse(chunk: dict) -> str:
    return f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"


def _synthetic_stream(text: str, pdf_url: str, model_id: str, title: str) -> StreamingResponse:
    """案内文を OpenAI 互換 SSE チャンクに変換する。

    process_chat_response (streaming_chat_response_handler) がこのチャンクを
    通常の LLM 応答と同じように処理するため、UI 配信・DB 保存・タイトル生成が
    標準経路で動く。PDF URL は delta.annotations (url_citation) で
    「ソース」カードとしても表示される。
    """

    async def gen():
        now = int(time.time())
        base = {"id": "chatcmpl-shineos-doc", "object": "chat.completion.chunk", "created": now, "model": model_id}
        delta = {"content": text}
        if pdf_url:
            fname = pdf_url.rsplit("/", 1)[-1]
            delta["annotations"] = [
                {
                    "type": "url_citation",
                    "url_citation": {"url": pdf_url, "title": f"{title}（PDF）" if title else fname},
                }
            ]
        yield _sse({**base, "choices": [{"index": 0, "delta": delta}]})
        yield _sse(
            {
                **base,
                "choices": [{"index": 0, "delta": {}, "finish_reason": "stop"}],
                "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
            }
        )
        yield "data: [DONE]\n\n"

    return StreamingResponse(gen(), media_type="text/event-stream")


def _synthetic_completion(text: str, model_id: str) -> dict:
    """非ストリーム (API) 用の OpenAI 互換レスポンス"""
    return {
        "id": "chatcmpl-shineos-doc",
        "object": "chat.completion",
        "created": int(time.time()),
        "model": model_id,
        "choices": [
            {
                "index": 0,
                "message": {"role": "assistant", "content": text},
                "finish_reason": "stop",
            }
        ],
        "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
    }


async def maybe_build_doc_response(request, form_data: dict, user, metadata: dict, model: dict):
    """process_chat の LLM 呼び出し前に呼ばれる入口。

    最後のユーザーメッセージが /doc で始まれば資料パイプラインを実行し、
    案内文の合成レスポンスを返す。/doc 以外は None を返す（通常チャット）。
    パイプライン失敗時はエラー案内の合成レスポンスを返す
    （Q&A 経路にフォールバックすると「対象外」拒否になってしまうため）。
    """
    # この関数は process_chat の先頭（payload 処理の前）で呼ばれるため、
    # messages / metadata['user_message'] には生の質問文が入っている。
    topic = ""
    um = metadata.get("user_message") or {}
    c = um.get("content") if isinstance(um, dict) else None
    if isinstance(c, str):
        topic = c.strip()
    elif isinstance(c, list):
        topic = " ".join(p.get("text", "") for p in c if isinstance(p, dict)).strip()
    if not topic:
        try:
            from open_webui.utils.misc import get_last_user_message

            topic = (get_last_user_message(form_data.get("messages", [])) or "").strip()
        except Exception:
            pass
    # トリガー判定: 長いコマンドから照合（/docx と /doc の前置混同を防ぐ）
    fmt = None
    for trigger in sorted(TRIGGERS, key=len, reverse=True):
        if topic.startswith(trigger):
            fmt = TRIGGERS[trigger]
            topic = topic[len(trigger):].strip()
            break
    if fmt is None:
        return None

    model_id = form_data.get("model") or ""
    emitter = None
    try:
        from open_webui.socket.main import get_event_emitter

        emitter = await get_event_emitter(metadata)
    except Exception:
        emitter = None

    try:
        if not topic:
            text = "主題が空のため資料を作成できませんでした。/pdf・/pptx・/docx の後に半角スペースと主題を付けて、もう一度送信してください。"
            return _finish(request, form_data, text, None, model_id, emitter, "")
        await _emit(emitter, "資料の目次を作成しています…")
        outline = await _make_outline(request, model_id, topic)
        sections = [{"heading": s["heading"], "query": s["query"], "body": ""} for s in outline["sections"]]

        for i, sec in enumerate(sections, 1):
            await _emit(emitter, f"第{i}節「{sec['heading']}」を執筆しています…")
            sec["body"] = await _write_section(request, form_data, sec["heading"], sec["query"])

        await _emit(emitter, f"{fmt.upper()} を生成しています…")
        url = await asyncio.to_thread(_build_document, fmt, outline["title"], sections)

        empty = [s["heading"] for s in sections if "該当記載がありません" in s["body"]]
        toc = "\n".join(f"{i}. {s['heading']}" for i, s in enumerate(sections, 1))
        note = ""
        if empty:
            note = "\n\n※ 「" + "」「".join(empty) + "」はナレッジに材料が見つからなかったため、該当記載なしとしています。"
        notice = (
            f"資料を作成しました（{fmt.upper()} / {len(sections)}セクション）。\n\n"
            f"**目次**\n{toc}\n\n"
            f"ファイルを開く: {url}\n"
            f"（クリックで開かない場合は URL をコピーしてブラウザで開いてください）"
            f"{note}"
        )
        await _emit(emitter, "資料を作成しました", done=True)
        return _finish(request, form_data, notice, url, model_id, emitter, outline["title"])
    except Exception:
        return _finish(
            request,
            form_data,
            "資料の作成中にエラーが発生しました。しばらく待ってからもう一度お試しください。",
            None,
            model_id,
            emitter,
            "",
        )


def _finish(request, form_data: dict, text: str, pdf_url, model_id: str, emitter, title: str):
    if form_data.get("stream"):
        return _synthetic_stream(text, pdf_url, model_id, title)
    return _synthetic_completion(text, model_id)
